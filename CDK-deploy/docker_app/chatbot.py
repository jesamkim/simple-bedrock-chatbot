import random
import json
from typing import List, Tuple, Union
import os
import tempfile
import csv
import pandas as pd

import streamlit as st
from langchain_community.chat_message_histories import StreamlitChatMessageHistory
from langchain_aws import ChatBedrock, ChatBedrockConverse
from langchain_core.messages import AIMessage, HumanMessage, SystemMessage
from langchain.callbacks.base import BaseCallbackHandler
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder
import boto3

REGION = "us-west-2"
MODELS = {
    "Claude 3.7 Sonnet": {
        "id": "us.anthropic.claude-3-7-sonnet-20250219-v1:0",
        "class": ChatBedrock,
        "use_model_kwargs": True,
        "max_tokens": 8192
    },
    "Nova Pro 1.0": {
        "id": "us.amazon.nova-pro-v1:0",
        "class": None,  # 직접 boto3 client 사용
        "use_model_kwargs": False,
        "max_tokens": 5000
    }
}

SUPPORTED_FORMATS = ['pdf', 'doc', 'docx', 'md', 'ppt', 'pptx', 'txt', 'html', 'csv', 'xls', 'xlsx']
MAX_MESSAGES = 10  # 대화 기록 최대 유지 수

class ChatMessage:
    def __init__(self, role: str, text: str):
        self.role = role
        self.text = text

def initialize_session_state():
    """세션 상태 초기화"""
    if "widget_key" not in st.session_state:
        st.session_state["widget_key"] = str(random.randint(1, 1000000))
    if "messages" not in st.session_state:
        st.session_state.messages = []
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = StreamlitChatMessageHistory(key="chat_history")
    if "context" not in st.session_state:
        st.session_state.context = None
    if "nova_messages" not in st.session_state:
        st.session_state.nova_messages = []
    if "initial_system_message" not in st.session_state:
        st.session_state.initial_system_message = None

class StreamHandler(BaseCallbackHandler):
    def __init__(self, container: st.container) -> None:
        self.container = container
        self.text = ""

    def on_llm_new_token(self, token: str, **kwargs) -> None:
        self.text += token
        self.container.markdown(self.text)

def set_page_config() -> None:
    st.set_page_config(page_title="Bedrock Chatbot", layout="wide")
    st.title("Bedrock Chatbot with Document Q&A")

def get_sidebar_params() -> Tuple[float, float, int, int, int, str]:
    with st.sidebar:
        st.markdown("## Model Selection")
        model_name = st.radio(
            "Choose a model",
            options=list(MODELS.keys()),
            index=0,
            key=f"{st.session_state['widget_key']}_Model"
        )
        
        st.markdown("## Document Upload")
        uploaded_file = st.file_uploader(
            "Upload a document",
            type=SUPPORTED_FORMATS,
            help="지원되는 파일 형식: pdf, doc, docx, md, ppt, pptx, txt, html, csv, xls, xlsx",
            key=f"{st.session_state['widget_key']}_file_uploader"  # widget_key와 연동
        )

        st.markdown("## Inference Parameters")
        system_prompt = st.text_area(
            "System Prompt",
            "You're a helpful assistant who loves to respond in Korean. You'll answer questions based on the uploaded documents, when a document exists. When you're not certain about something, don't guess or make up information. Instead, honestly admit that you don't know.",
            key=f"{st.session_state['widget_key']}_System_Prompt",
        )

        temperature = st.slider(
            "Temperature",
            min_value=0.0,
            max_value=1.0,
            value=0.0,
            step=0.1,
            key=f"{st.session_state['widget_key']}_Temperature",
        )

        with st.container():
            col1, col2 = st.columns(2)
            with col1:
                top_p = st.slider(
                    "Top-P",
                    min_value=0.0,
                    max_value=1.0,
                    value=1.00,
                    step=0.01,
                    key=f"{st.session_state['widget_key']}_Top-P",
                )
            with col2:
                top_k = st.slider(
                    "Top-K",
                    min_value=1,
                    max_value=500,
                    value=500,
                    step=5,
                    key=f"{st.session_state['widget_key']}_Top-K",
                )

        with st.container():
            col1, col2 = st.columns(2)
            with col1:
                model_max_tokens = MODELS[model_name]["max_tokens"]
                max_tokens = st.slider(
                    "Max Token",
                    min_value=0,
                    max_value=model_max_tokens,
                    value=model_max_tokens,
                    step=8,
                    key=f"{st.session_state['widget_key']}_Max_Token",
                )
            with col2:
                memory_window = st.slider(
                    "Memory Window",
                    min_value=0,
                    max_value=10,
                    value=10,
                    step=1,
                    key=f"{st.session_state['widget_key']}_Memory_Window",
                )

    return temperature, top_p, top_k, max_tokens, memory_window, system_prompt, uploaded_file, model_name

def process_uploaded_file(file_path: str) -> str:
    """문서 파일을 처리하여 텍스트로 변환"""
    file_ext = os.path.splitext(file_path)[1].lower()
    
    try:
        if file_ext == '.pdf':
            from PyPDF2 import PdfReader
            reader = PdfReader(file_path)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            return text
        
        elif file_ext in ['.doc', '.docx']:
            from docx import Document
            doc = Document(file_path)
            text = ""
            for para in doc.paragraphs:
                text += para.text + "\n"
            return text
        
        elif file_ext in ['.txt', '.md', '.html']:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        
        elif file_ext in ['.csv']:
            df = pd.read_csv(file_path)
            return df.to_string()
        
        elif file_ext in ['.xls', '.xlsx']:
            df = pd.read_excel(file_path)
            return df.to_string()
        
        elif file_ext in ['.ppt', '.pptx']:
            from pptx import Presentation
            text = []
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n\n".join(text)
        
        else:
            raise ValueError(f"지원하지 않는 파일 형식입니다: {file_ext}")
            
    except Exception as e:
        raise Exception(f"파일 처리 중 오류가 발생했습니다 ({file_ext}): {str(e)}")

def init_conversation_chain(
    temperature: float,
    top_p: float,
    top_k: int,
    max_tokens: int,
    system_prompt: str,
    model_name: str,
) -> Union[ChatBedrock, boto3.client]:

    model_info = MODELS[model_name]
    model_id = model_info["id"]
    use_model_kwargs = model_info["use_model_kwargs"]

    if use_model_kwargs:  # Claude 3.7 Sonnet
        model_kwargs = {
            "temperature": temperature,
            "top_p": top_p,
            "top_k": top_k,
            "max_tokens": max_tokens,
            "system": system_prompt
        }
        return ChatBedrock(
            model_id=model_id,
            model_kwargs=model_kwargs,
            streaming=True,
            region_name=REGION
        )
    else:  # Nova Pro
        return boto3.client("bedrock-runtime", region_name=REGION)

def convert_chat_messages_to_converse_api(chat_messages: List[ChatMessage]) -> List[dict]:
    """ChatMessage 객체 리스트를 Nova Pro API 형식으로 변환"""
    messages = []
    for chat_msg in chat_messages:
        if chat_msg.role != 'system':  # system 역할은 건너뛰기
            messages.append({
                "role": chat_msg.role,
                "content": [
                    {
                        "text": chat_msg.text
                    }
                ]
            })
    return messages

def generate_response(
    conversation: Union[ChatBedrock, boto3.client],
    input_text: str,
    chat_history: StreamlitChatMessageHistory
) -> str:
    with st.chat_message("assistant"):
        message_placeholder = st.empty()
        full_response = ""
        try:
            if isinstance(conversation, ChatBedrock):  # Claude 3.7 Sonnet
                messages = []
                messages.append(SystemMessage(content=conversation.model_kwargs["system"]))
                for msg in chat_history.messages:
                    messages.append(msg)
                messages.append(HumanMessage(content=input_text))
                
                for chunk in conversation.stream(messages):
                    if chunk.content:
                        full_response += chunk.content
                        message_placeholder.markdown(full_response + "▌")
                message_placeholder.markdown(full_response)
                return full_response
            else:  # Nova Pro
                # 새로운 사용자 메시지 추가
                new_message = ChatMessage('user', input_text)
                st.session_state.nova_messages.append(new_message)
                
                # 메시지 수 제한 확인
                if len(st.session_state.nova_messages) > MAX_MESSAGES * 2:  # user와 assistant 메시지 쌍을 고려
                    st.session_state.nova_messages = st.session_state.nova_messages[-(MAX_MESSAGES * 2):]
                
                # Nova Pro API 형식으로 메시지 변환
                messages = convert_chat_messages_to_converse_api(st.session_state.nova_messages)
                
                # Nova Pro 파라미터 설정
                widget_key = st.session_state["widget_key"]
                
                response = conversation.converse_stream(
                    modelId=MODELS["Nova Pro 1.0"]["id"],
                    messages=messages,
                    inferenceConfig={
                        "temperature": float(st.session_state[f"{widget_key}_Temperature"]),
                        "maxTokens": int(st.session_state[f"{widget_key}_Max_Token"]),
                    }
                )

                # 스트리밍 응답 처리
                for event in response['stream']:
                    if 'contentBlockDelta' in event:
                        chunk = event['contentBlockDelta']['delta']['text']
                        full_response += chunk
                        message_placeholder.markdown(full_response + "▌")
                
                message_placeholder.markdown(full_response)
                
                # 어시스턴트 응답 저장
                st.session_state.nova_messages.append(ChatMessage('assistant', full_response))
                return full_response

        except Exception as e:
            if "ThrottlingException" in str(e):
                error_message = "요청을 처리하지 못했습니다. 잠시 후 다시 말씀해 주세요. 🙏"
            else:
                error_message = f"죄송합니다. 오류가 발생했습니다: {str(e)}"
            message_placeholder.markdown(error_message)
            return error_message

def new_chat() -> None:
    """새로운 대화 시작"""
    # widget_key를 새로 생성하여 file_uploader를 포함한 모든 위젯 초기화
    st.session_state["widget_key"] = str(random.randint(1, 1000000))
    st.session_state.messages = []
    st.session_state.chat_history.clear()
    st.session_state.nova_messages = []
    st.session_state.context = None
    st.session_state.initial_system_message = None

def handle_file_upload(uploaded_file) -> str:
    """파일 업로드 처리"""
    if uploaded_file:
        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(uploaded_file.name)[1]) as tmp_file:
            tmp_file.write(uploaded_file.getvalue())
            tmp_file_path = tmp_file.name

        try:
            document_text = process_uploaded_file(tmp_file_path)
            st.session_state.context = document_text
            os.unlink(tmp_file_path)
            return document_text
        except Exception as e:
            os.unlink(tmp_file_path)
            st.error(f"문서 처리 중 오류가 발생했습니다: {str(e)}")
            return None
    return None

def main() -> None:
    set_page_config()
    initialize_session_state()

    st.sidebar.button("New Chat", on_click=new_chat, type="primary")

    temperature, top_p, top_k, max_tokens, memory_window, system_prompt, uploaded_file, model_name = get_sidebar_params()

    # 문서가 업로드되면 시스템 메시지 초기화
    if uploaded_file:
        document_context = handle_file_upload(uploaded_file)
        if document_context:
            st.sidebar.success(f"문서가 성공적으로 업로드되었습니다: {uploaded_file.name}")
            # 문서 컨텍스트를 포함한 시스템 메시지 생성
            full_system_prompt = f"{system_prompt}\n\n참고할 문서 내용:\n\n{document_context}"
            st.session_state.initial_system_message = full_system_prompt
            
            # Nova Pro의 경우 초기 시스템 메시지를 user 메시지로 추가
            if model_name == "Nova Pro 1.0" and not st.session_state.nova_messages:
                st.session_state.nova_messages = [
                    ChatMessage('user', full_system_prompt),
                    ChatMessage('assistant', "네, 이해했습니다. 업로드된 문서를 참고하여 답변하도록 하겠습니다.")
                ]

    # 시스템 프롬프트 설정
    if st.session_state.initial_system_message and model_name == "Claude 3.7 Sonnet":
        system_prompt = st.session_state.initial_system_message

    conv_chain = init_conversation_chain(
        temperature, top_p, top_k, max_tokens, system_prompt, model_name
    )

    # 저장된 메시지 표시
    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])

    # 사용자 입력 처리
    if prompt := st.chat_input("메시지를 입력하세요"):
        st.session_state.messages.append({"role": "user", "content": prompt})
        st.session_state.chat_history.add_user_message(prompt)
        
        with st.chat_message("user"):
            st.markdown(prompt)

        # 응답 생성
        response = generate_response(conv_chain, prompt, st.session_state.chat_history)
        
        st.session_state.messages.append({"role": "assistant", "content": response})
        st.session_state.chat_history.add_ai_message(response)

if __name__ == "__main__":
    main()
