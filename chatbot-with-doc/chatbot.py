import random
import json
from typing import List, Tuple, Union
import os
import tempfile
import csv
import pandas as pd
import warnings

# Pydantic 경고 무시
warnings.filterwarnings("ignore", category=UserWarning, module="pydantic")

import streamlit as st
from langchain_community.chat_message_histories import StreamlitChatMessageHistory
from langchain_aws import ChatBedrock, ChatBedrockConverse
from langchain_core.messages import AIMessage, HumanMessage, SystemMessage
from langchain.callbacks.base import BaseCallbackHandler
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder
import boto3

REGION = "us-west-2"
MODELS = {
    "Claude 3.7 Sonnet": {
        "id": "us.anthropic.claude-3-7-sonnet-20250219-v1:0",
        "class": ChatBedrock,
        "use_model_kwargs": True,
        "max_tokens": 8192
    },
    "Nova Pro 1.0": {
        "id": "us.amazon.nova-pro-v1:0",
        "class": None,  # 직접 boto3 client 사용
        "use_model_kwargs": False,
        "max_tokens": 5000
    }
}

SUPPORTED_FORMATS = ['pdf', 'doc', 'docx', 'md', 'ppt', 'pptx', 'txt', 'html', 'csv', 'xls', 'xlsx']
MAX_MESSAGES = 10  # 대화 기록 최대 유지 수

class ChatMessage:
    def __init__(self, role: str, text: str):
        self.role = role
        self.text = text

def initialize_session_state():
    """세션 상태 초기화"""
    if "widget_key" not in st.session_state:
        st.session_state["widget_key"] = str(random.randint(1, 1000000))
    if "messages" not in st.session_state:
        st.session_state.messages = []
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = StreamlitChatMessageHistory(key="chat_history")
    if "context" not in st.session_state:
        st.session_state.context = None
    if "nova_messages" not in st.session_state:
        st.session_state.nova_messages = []
    if "initial_system_message" not in st.session_state:
        st.session_state.initial_system_message = None

class StreamHandler(BaseCallbackHandler):
    def __init__(self, container: st.container) -> None:
        self.container = container
        self.text = ""

    def on_llm_new_token(self, token: str, **kwargs) -> None:
        self.text += token
        self.container.markdown(self.text)

def set_page_config() -> None:
    st.set_page_config(page_title="Bedrock Chatbot", layout="wide")
    st.title("Bedrock Chatbot with Document Q&A")

def get_sidebar_params() -> Tuple[float, float, int, int, int, str, object, str, bool, bool]:
    with st.sidebar:
        st.markdown("## Model Selection")
        model_name = st.radio(
            "Choose a model",
            options=list(MODELS.keys()),
            index=0,
            key=f"{st.session_state['widget_key']}_Model"
        )
        
        # Claude 3.7 Sonnet 모델이 선택된 경우에만 Model reasoning 모드 옵션 표시
        extended_thinking = False  # 기본값으로 비활성화
        show_reasoning = False  # 기본값으로 비활성화
        if model_name == "Claude 3.7 Sonnet":
            extended_thinking = st.checkbox(
                "Model reasoning 모드 활성화",
                value=False,  # 기본값으로 비활성화
                help="Claude 3.7 Sonnet의 추론 모드를 활성화합니다. 복잡한 문제 해결에 도움이 됩니다. (참고: 이 모드에서는 Temperature가 자동으로 1.0으로 설정되고 Top-K 및 Top-P 설정이 비활성화됩니다)",
                key=f"{st.session_state['widget_key']}_Model_Reasoning"
            )
            
            if extended_thinking:
                st.info("Model reasoning 모드가 활성화되어 Temperature 값이 1.0으로 자동 설정되고 Top-K 및 Top-P 설정이 비활성화됩니다.")
                
                show_reasoning = st.checkbox(
                    "Reasoning 과정 표시",
                    value=True,
                    help="Claude의 사고 과정을 실시간으로 표시합니다.",
                    key=f"{st.session_state['widget_key']}_Show_Reasoning"
                )
        
        st.markdown("## Document Upload")
        uploaded_file = st.file_uploader(
            "Upload a document",
            type=SUPPORTED_FORMATS,
            help="지원되는 파일 형식: pdf, doc, docx, md, ppt, pptx, txt, html, csv, xls, xlsx",
            key=f"{st.session_state['widget_key']}_file_uploader"  # widget_key와 연동
        )

        st.markdown("## Inference Parameters")
        system_prompt = st.text_area(
            "System Prompt",
            "You're a helpful assistant who loves to respond in Korean. You'll answer questions based on the uploaded documents, when a document exists. When you're not certain about something, don't guess or make up information. Instead, honestly admit that you don't know.",
            key=f"{st.session_state['widget_key']}_System_Prompt",
        )

        temperature = st.slider(
            "Temperature",
            min_value=0.0,
            max_value=1.0,
            value=0.0,
            step=0.1,
            key=f"{st.session_state['widget_key']}_Temperature",
        )

        with st.container():
            col1, col2 = st.columns(2)
            with col1:
                top_p = st.slider(
                    "Top-P",
                    min_value=0.0,
                    max_value=1.0,
                    value=1.00,
                    step=0.01,
                    key=f"{st.session_state['widget_key']}_Top-P",
                )
            with col2:
                top_k = st.slider(
                    "Top-K",
                    min_value=1,
                    max_value=500,
                    value=500,
                    step=5,
                    key=f"{st.session_state['widget_key']}_Top-K",
                )

        with st.container():
            col1, col2 = st.columns(2)
            with col1:
                model_max_tokens = MODELS[model_name]["max_tokens"]
                max_tokens = st.slider(
                    "Max Token",
                    min_value=0,
                    max_value=model_max_tokens,
                    value=model_max_tokens,
                    step=8,
                    key=f"{st.session_state['widget_key']}_Max_Token",
                )
            with col2:
                memory_window = st.slider(
                    "Memory Window",
                    min_value=0,
                    max_value=10,
                    value=10,
                    step=1,
                    key=f"{st.session_state['widget_key']}_Memory_Window",
                )

    return temperature, top_p, top_k, max_tokens, memory_window, system_prompt, uploaded_file, model_name, extended_thinking, show_reasoning

def process_uploaded_file(file_path: str) -> str:
    """문서 파일을 처리하여 텍스트로 변환"""
    file_ext = os.path.splitext(file_path)[1].lower()
    
    try:
        if file_ext == '.pdf':
            from PyPDF2 import PdfReader
            reader = PdfReader(file_path)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            return text
        
        elif file_ext in ['.doc', '.docx']:
            from docx import Document
            doc = Document(file_path)
            text = ""
            for para in doc.paragraphs:
                text += para.text + "\n"
            return text
        
        elif file_ext in ['.txt', '.md', '.html']:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        
        elif file_ext in ['.csv']:
            df = pd.read_csv(file_path)
            return df.to_string()
        
        elif file_ext in ['.xls', '.xlsx']:
            df = pd.read_excel(file_path)
            return df.to_string()
        
        elif file_ext in ['.ppt', '.pptx']:
            from pptx import Presentation
            text = []
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n\n".join(text)
        
        else:
            raise ValueError(f"지원하지 않는 파일 형식입니다: {file_ext}")
            
    except Exception as e:
        raise Exception(f"파일 처리 중 오류가 발생했습니다 ({file_ext}): {str(e)}")

def init_conversation_chain(
    temperature: float,
    top_p: float,
    top_k: int,
    max_tokens: int,
    system_prompt: str,
    model_name: str,
    extended_thinking: bool = False,
) -> Union[ChatBedrock, boto3.client]:

    model_info = MODELS[model_name]
    model_id = model_info["id"]
    use_model_kwargs = model_info["use_model_kwargs"]

    if use_model_kwargs:  # Claude 3.7 Sonnet
        model_kwargs = {
            "temperature": temperature,
            "top_p": top_p,
            "top_k": top_k,
            "max_tokens": max_tokens,
            "system": system_prompt
        }
        
        # Model reasoning 모드가 활성화된 경우 (Claude 3.7 Sonnet 전용)
        if model_name == "Claude 3.7 Sonnet" and extended_thinking:
            # Model reasoning 모드에서는 temperature가 반드시 1이어야 함
            model_kwargs["temperature"] = 1.0
            model_kwargs["anthropic_version"] = "bedrock-2023-05-31"
            
            # Model reasoning 모드에서는 top_k와 top_p를 설정하지 않아야 함
            if "top_k" in model_kwargs:
                del model_kwargs["top_k"]
            if "top_p" in model_kwargs:
                del model_kwargs["top_p"]
            
            # 최대 Length는 64000으로 설정
            model_kwargs["max_tokens"] = 64000
            
            # thinking.budget_tokens는 max_tokens보다 작아야 하며, 최대 4096
            thinking_budget = min(4096, model_kwargs["max_tokens"] - 1000)
            
            model_kwargs["thinking"] = {
                "type": "enabled",
                "budget_tokens": thinking_budget  # 사고 과정에 할당할 토큰 수
            }
            
        return ChatBedrock(
            model_id=model_id,
            model_kwargs=model_kwargs,
            streaming=True,
            region_name=REGION
        )
    else:  # Nova Pro
        return boto3.client("bedrock-runtime", region_name=REGION)

def convert_langchain_messages_to_anthropic(messages):
    """LangChain 메시지를 Anthropic API 형식으로 변환"""
    result = []
    system_message = None
    
    for msg in messages:
        if msg.type == "system":
            system_message = msg.content  # 시스템 메시지 별도 저장
        elif msg.type == "ai":
            result.append({"role": "assistant", "content": msg.content})
        elif msg.type == "human":
            result.append({"role": "user", "content": msg.content})
    
    return result, system_message

def convert_chat_messages_to_converse_api(chat_messages: List[ChatMessage]) -> List[dict]:
    """ChatMessage 객체 리스트를 Nova Pro API 형식으로 변환"""
    messages = []
    for chat_msg in chat_messages:
        if chat_msg.role != 'system':  # system 역할은 건너뛰기
            messages.append({
                "role": chat_msg.role,
                "content": [
                    {
                        "text": chat_msg.text
                    }
                ]
            })
    return messages

def generate_response(
    conversation: Union[ChatBedrock, boto3.client],
    input_text: str,
    chat_history: StreamlitChatMessageHistory,
    show_reasoning: bool = False
) -> str:
    with st.chat_message("assistant"):
        message_placeholder = st.empty()
        full_response = ""
        try:
            if isinstance(conversation, ChatBedrock):  # Claude 3.7 Sonnet
                messages = []
                messages.append(SystemMessage(content=conversation.model_kwargs["system"]))
                for msg in chat_history.messages:
                    messages.append(msg)
                messages.append(HumanMessage(content=input_text))
                
                # Model reasoning 모드 활성화 여부 확인
                has_thinking = "thinking" in conversation.model_kwargs
                
                try:
                    # 스트리밍 처리 수정
                    if has_thinking:
                        # Model reasoning 모드에서는 직접 boto3 클라이언트 사용
                        bedrock_runtime = boto3.client("bedrock-runtime", region_name=REGION)
                        
                        # 메시지 변환
                        anthropic_messages, system_content = convert_langchain_messages_to_anthropic(messages)
                        
                        # 최대 토큰 값 설정 (64000)
                        model_max_tokens = 64000
                        
                        request_payload = {
                            "anthropic_version": "bedrock-2023-05-31",
                            "max_tokens": model_max_tokens,
                            "temperature": 1.0,
                            "system": system_content,  # 시스템 메시지를 별도 파라미터로 전달
                            "thinking": conversation.model_kwargs["thinking"],
                            "messages": anthropic_messages
                        }
                        
                        response = bedrock_runtime.invoke_model_with_response_stream(
                            modelId=MODELS["Claude 3.7 Sonnet"]["id"],
                            body=json.dumps(request_payload)
                        )
                        
                        # 디버깅을 위한 요청 페이로드 출력
                        print(f"요청 페이로드: {json.dumps(request_payload, indent=2)}")
                        
                        # 채팅 메시지 위에 reasoning을 표시하기 위해 reasoning_placeholder를 먼저 생성
                        reasoning_placeholder = st.empty()
                        
                        reasoning_text = ""
                        has_shown_reasoning = False
                        
                        # 디버깅을 위한 로그 추가
                        print(f"show_reasoning 값: {show_reasoning}")
                        
                        for event in response["body"]:
                            try:
                                chunk = json.loads(event["chunk"]["bytes"])
                                print(f"청크 데이터: {chunk}")  # 디버깅용
                                
                                # thinking 타입 처리 (reasoning 과정)
                                if chunk.get("type") == "thinking":
                                    thinking_content = chunk.get("thinking", "")
                                    reasoning_text += thinking_content
                                    print(f"사고 과정 감지: {thinking_content}")
                                    has_shown_reasoning = True
                                    
                                    # reasoning 과정을 UI에 표시 (show_reasoning이 True인 경우에만)
                                    if show_reasoning:
                                        # 메시지 플레이스홀더 위에 reasoning 표시
                                        # 라이트/다크 테마에 따라 다른 색상 사용
                                        is_dark_theme = True  # 기본값으로 다크 테마 가정
                                        try:
                                            # Streamlit 테마 감지 시도
                                            theme = st.get_option("theme.base")
                                            is_dark_theme = theme == "dark"
                                        except:
                                            pass  # 오류 발생 시 기본값 사용
                                        
                                        bg_color = "#2a2a2a" if is_dark_theme else "#f0f2f6"
                                        text_color = "#ffffff" if is_dark_theme else "#262730"
                                        
                                        reasoning_placeholder.markdown(f"""
                                        <div style="background-color: {bg_color}; color: {text_color}; padding: 10px; border-radius: 5px; margin-bottom: 10px; border-left: 4px solid #007bff; border: 1px solid #007bff; max-height: 400px; overflow-y: auto;">
                                            <h4 style="color: {text_color}; margin-top: 0;">🧠 Reasoning...</h4>
                                            <pre style="white-space: pre-wrap; overflow-wrap: break-word; color: {text_color}; background-color: transparent; border: none; padding: 0; margin: 0;">{reasoning_text}</pre>
                                        </div>
                                        """, unsafe_allow_html=True)
                                        print("Reasoning UI 업데이트됨")
                                
                                # content_block_delta 타입 처리 - text 또는 thinking_delta 모두 처리
                                elif chunk.get("type") == "content_block_delta":
                                    # text 또는 text_delta 처리
                                    if chunk["delta"].get("type") == "text" or chunk["delta"].get("type") == "text_delta":
                                        text_chunk = chunk["delta"].get("text", "")
                                        full_response += text_chunk
                                        print(f"현재 응답: {full_response}")  # 응답 내용 확인
                                        
                                        # 매 청크마다 업데이트하지 말고 일정 간격으로 업데이트
                                        if len(text_chunk) > 10 or text_chunk.endswith(('.', '!', '?', '\n')):
                                            message_placeholder.markdown(full_response + "▌")
                                    
                                    # thinking_delta 처리 (사고 과정)
                                    elif show_reasoning and chunk["delta"].get("type") == "thinking_delta":
                                        thinking_chunk = chunk["delta"].get("thinking", "")
                                        if thinking_chunk:
                                            reasoning_text += thinking_chunk
                                            print(f"thinking_delta 사고 과정: {thinking_chunk}")
                                            has_shown_reasoning = True
                                            
                                            # reasoning 과정을 UI에 표시
                                            # 라이트/다크 테마에 따라 다른 색상 사용
                                            is_dark_theme = True  # 기본값으로 다크 테마 가정
                                            try:
                                                # Streamlit 테마 감지 시도
                                                theme = st.get_option("theme.base")
                                                is_dark_theme = theme == "dark"
                                            except:
                                                pass  # 오류 발생 시 기본값 사용
                                            
                                            bg_color = "#2a2a2a" if is_dark_theme else "#f0f2f6"
                                            text_color = "#ffffff" if is_dark_theme else "#262730"
                                            
                                            reasoning_placeholder.markdown(f"""
                                            <div style="background-color: {bg_color}; color: {text_color}; padding: 10px; border-radius: 5px; margin-bottom: 10px; border-left: 4px solid #007bff; border: 1px solid #007bff; max-height: 400px; overflow-y: auto;">
                                                <h4 style="color: {text_color}; margin-top: 0;">🧠 Reasoning...</h4>
                                                <pre style="white-space: pre-wrap; overflow-wrap: break-word; color: {text_color}; background-color: transparent; border: none; padding: 0; margin: 0;">{reasoning_text}</pre>
                                            </div>
                                            """, unsafe_allow_html=True)
                                
                                # 다른 형식의 thinking 데이터 처리 시도
                                elif "thinking" in chunk:
                                    thinking_content = chunk.get("thinking", "")
                                    reasoning_text += thinking_content
                                    print(f"다른 형식의 사고 과정 감지: {thinking_content}")
                                    has_shown_reasoning = True
                                    
                                    # reasoning 과정을 UI에 표시 (show_reasoning이 True인 경우에만)
                                    if show_reasoning:
                                        # 라이트/다크 테마에 따라 다른 색상 사용
                                        is_dark_theme = True  # 기본값으로 다크 테마 가정
                                        try:
                                            # Streamlit 테마 감지 시도
                                            theme = st.get_option("theme.base")
                                            is_dark_theme = theme == "dark"
                                        except:
                                            pass  # 오류 발생 시 기본값 사용
                                        
                                        bg_color = "#2a2a2a" if is_dark_theme else "#f0f2f6"
                                        text_color = "#ffffff" if is_dark_theme else "#262730"
                                        
                                        reasoning_placeholder.markdown(f"""
                                        <div style="background-color: {bg_color}; color: {text_color}; padding: 10px; border-radius: 5px; margin-bottom: 10px; border-left: 4px solid #007bff; border: 1px solid #007bff; max-height: 400px; overflow-y: auto;">
                                            <h4 style="color: {text_color}; margin-top: 0;">🧠 Reasoning...</h4>
                                            <pre style="white-space: pre-wrap; overflow-wrap: break-word; color: {text_color}; background-color: transparent; border: none; padding: 0; margin: 0;">{reasoning_text}</pre>
                                        </div>
                                        """, unsafe_allow_html=True)
                                        print("Reasoning UI 업데이트됨 (다른 형식)")
                                
                                # content_block_start 타입 처리 (reasoning 과정이 여기에 포함될 수 있음)
                                elif chunk.get("type") == "content_block_start" and "thinking" in str(chunk):
                                    try:
                                        # 다양한 형식의 thinking 데이터 추출 시도
                                        thinking_content = ""
                                        
                                        # content_block 내에 thinking 필드가 있는 경우
                                        if "content_block" in chunk and "thinking" in chunk["content_block"]:
                                            thinking_content = chunk["content_block"].get("thinking", "")
                                        # 직접 thinking 필드가 있는 경우
                                        elif "thinking" in chunk:
                                            thinking_content = chunk.get("thinking", "")
                                        
                                        # 추출된 내용이 있는 경우에만 추가
                                        if thinking_content and not isinstance(thinking_content, dict):
                                            reasoning_text += thinking_content
                                            print(f"content_block_start에서 사고 과정 감지: {thinking_content}")
                                            has_shown_reasoning = True
                                            
                                            # reasoning 과정을 UI에 표시 (show_reasoning이 True인 경우에만)
                                            if show_reasoning:
                                                # 라이트/다크 테마에 따라 다른 색상 사용
                                                is_dark_theme = True  # 기본값으로 다크 테마 가정
                                                try:
                                                    # Streamlit 테마 감지 시도
                                                    theme = st.get_option("theme.base")
                                                    is_dark_theme = theme == "dark"
                                                except:
                                                    pass  # 오류 발생 시 기본값 사용
                                                
                                                bg_color = "#2a2a2a" if is_dark_theme else "#f0f2f6"
                                                text_color = "#ffffff" if is_dark_theme else "#262730"
                                                
                                                reasoning_placeholder.markdown(f"""
                                                <div style="background-color: {bg_color}; color: {text_color}; padding: 10px; border-radius: 5px; margin-bottom: 10px; border-left: 4px solid #007bff; border: 1px solid #007bff; max-height: 400px; overflow-y: auto;">
                                                    <h4 style="color: {text_color}; margin-top: 0;">🧠 Reasoning...</h4>
                                                    <pre style="white-space: pre-wrap; overflow-wrap: break-word; color: {text_color}; background-color: transparent; border: none; padding: 0; margin: 0;">{reasoning_text}</pre>
                                                </div>
                                                """, unsafe_allow_html=True)
                                                print("Reasoning UI 업데이트됨 (content_block_start)")
                                    except Exception as thinking_error:
                                        print(f"사고 과정 추출 오류: {str(thinking_error)}")
                                
                                # 모든 청크에서 "thinking" 문자열을 찾아 처리 (마지막 시도)
                                elif show_reasoning and "thinking" in str(chunk).lower():
                                    try:
                                        # JSON 데이터에서 실제 사고 과정 내용만 추출
                                        chunk_data = chunk
                                        thinking_content = ""
                                        
                                        # content_block_delta 타입인 경우
                                        if isinstance(chunk_data, dict) and chunk_data.get("type") == "content_block_delta":
                                            if "delta" in chunk_data and "thinking" in chunk_data["delta"]:
                                                thinking_content = chunk_data["delta"].get("thinking", "")
                                        
                                        # 다른 형식의 thinking 데이터 처리
                                        elif isinstance(chunk_data, dict) and "thinking" in chunk_data:
                                            thinking_content = chunk_data.get("thinking", "")
                                        
                                        # 추출된 내용이 있는 경우에만 추가
                                        if thinking_content:
                                            reasoning_text += thinking_content
                                            print(f"사고 과정 내용 추출: {thinking_content}")
                                            has_shown_reasoning = True
                                            
                                            # reasoning 과정을 UI에 표시
                                            # 라이트/다크 테마에 따라 다른 색상 사용
                                            is_dark_theme = True  # 기본값으로 다크 테마 가정
                                            try:
                                                # Streamlit 테마 감지 시도
                                                theme = st.get_option("theme.base")
                                                is_dark_theme = theme == "dark"
                                            except:
                                                pass  # 오류 발생 시 기본값 사용
                                            
                                            bg_color = "#2a2a2a" if is_dark_theme else "#f0f2f6"
                                            text_color = "#ffffff" if is_dark_theme else "#262730"
                                            
                                            reasoning_placeholder.markdown(f"""
                                            <div style="background-color: {bg_color}; color: {text_color}; padding: 10px; border-radius: 5px; margin-bottom: 10px; border-left: 4px solid #007bff; border: 1px solid #007bff; max-height: 400px; overflow-y: auto;">
                                                <h4 style="color: {text_color}; margin-top: 0;">🧠 Reasoning...</h4>
                                                <pre style="white-space: pre-wrap; overflow-wrap: break-word; color: {text_color}; background-color: transparent; border: none; padding: 0; margin: 0;">{reasoning_text}</pre>
                                            </div>
                                            """, unsafe_allow_html=True)
                                            print("Reasoning UI 업데이트됨 (내용 추출)")
                                    except Exception as thinking_error:
                                        print(f"사고 과정 내용 추출 오류: {str(thinking_error)}")
                                        
                                # 응답 완료 이벤트 처리
                                elif chunk.get("type") == "message_stop":
                                    print("응답 완료")
                                    message_placeholder.markdown(full_response)
                            except Exception as chunk_error:
                                print(f"청크 처리 오류: {str(chunk_error)}")
                        
                        # 스트리밍 완료 후 최종 메시지 표시
                        print(f"최종 응답: {full_response}")
                        message_placeholder.markdown(full_response)
                        
                        # 디버깅 정보 출력
                        print(f"reasoning 표시 여부: {has_shown_reasoning}, show_reasoning 값: {show_reasoning}")
                        if not has_shown_reasoning and show_reasoning:
                            print("경고: reasoning 데이터가 감지되지 않았습니다.")
                    else:
                        # 일반 모드에서는 LangChain 스트리밍 사용
                        for chunk in conversation.stream(messages):
                            if hasattr(chunk, 'content') and chunk.content:
                                full_response += chunk.content
                                message_placeholder.markdown(full_response + "▌")
                    
                    message_placeholder.markdown(full_response)
                    return full_response
                except Exception as e:
                    st.error(f"스트리밍 처리 중 오류: {str(e)}")
                    # 스트리밍 실패 시 일반 호출 시도
                    try:
                        response = conversation.invoke(messages)
                        full_response = response.content
                        message_placeholder.markdown(full_response)
                        return full_response
                    except Exception as e2:
                        st.error(f"일반 호출 처리 중 오류: {str(e2)}")
                        return f"응답 생성 중 오류가 발생했습니다: {str(e2)}"
            else:  # Nova Pro
                # 새로운 사용자 메시지 추가
                new_message = ChatMessage('user', input_text)
                st.session_state.nova_messages.append(new_message)
                
                # 메시지 수 제한 확인
                if len(st.session_state.nova_messages) > MAX_MESSAGES * 2:  # user와 assistant 메시지 쌍을 고려
                    st.session_state.nova_messages = st.session_state.nova_messages[-(MAX_MESSAGES * 2):]
                
                # Nova Pro API 형식으로 메시지 변환
                messages = convert_chat_messages_to_converse_api(st.session_state.nova_messages)
                
                # Nova Pro 파라미터 설정
                widget_key = st.session_state["widget_key"]
                
                response = conversation.converse_stream(
                    modelId=MODELS["Nova Pro 1.0"]["id"],
                    messages=messages,
                    inferenceConfig={
                        "temperature": float(st.session_state[f"{widget_key}_Temperature"]),
                        "maxTokens": int(st.session_state[f"{widget_key}_Max_Token"]),
                    }
                )

                # 스트리밍 응답 처리
                for event in response['stream']:
                    if 'contentBlockDelta' in event:
                        chunk = event['contentBlockDelta']['delta']['text']
                        full_response += chunk
                        message_placeholder.markdown(full_response + "▌")
                
                message_placeholder.markdown(full_response)
                
                # 어시스턴트 응답 저장
                st.session_state.nova_messages.append(ChatMessage('assistant', full_response))
                return full_response

        except Exception as e:
            error_detail = str(e)
            print(f"상세 오류: {error_detail}")  # 터미널에 자세한 오류 출력
            
            if "ThrottlingException" in error_detail:
                error_message = "요청을 처리하지 못했습니다. 잠시 후 다시 말씀해 주세요. 🙏"
            elif "'message'" in error_detail:
                error_message = "응답 형식 오류가 발생했습니다. Model reasoning 모드 설정을 확인해주세요."
            elif "validationException" in error_detail:
                error_message = "API 검증 오류가 발생했습니다. 요청 형식을 확인해주세요."
                print(f"API 검증 오류 상세: {error_detail}")
            else:
                error_message = f"죄송합니다. 오류가 발생했습니다: {error_detail}"
            
            message_placeholder.markdown(error_message)
            return error_message

def new_chat() -> None:
    """새로운 대화 시작"""
    # widget_key를 새로 생성하여 file_uploader를 포함한 모든 위젯 초기화
    st.session_state["widget_key"] = str(random.randint(1, 1000000))
    st.session_state.messages = []
    st.session_state.chat_history.clear()
    st.session_state.nova_messages = []
    st.session_state.context = None
    st.session_state.initial_system_message = None

def handle_file_upload(uploaded_file) -> str:
    """파일 업로드 처리"""
    if uploaded_file:
        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(uploaded_file.name)[1]) as tmp_file:
            tmp_file.write(uploaded_file.getvalue())
            tmp_file_path = tmp_file.name

        try:
            document_text = process_uploaded_file(tmp_file_path)
            st.session_state.context = document_text
            os.unlink(tmp_file_path)
            return document_text
        except Exception as e:
            os.unlink(tmp_file_path)
            st.error(f"문서 처리 중 오류가 발생했습니다: {str(e)}")
            return None
    return None

def main() -> None:
    set_page_config()
    initialize_session_state()

    st.sidebar.button("New Chat", on_click=new_chat, type="primary")

    temperature, top_p, top_k, max_tokens, memory_window, system_prompt, uploaded_file, model_name, extended_thinking, show_reasoning = get_sidebar_params()

    # 문서가 업로드되면 시스템 메시지 초기화
    if uploaded_file:
        document_context = handle_file_upload(uploaded_file)
        if document_context:
            st.sidebar.success(f"문서가 성공적으로 업로드되었습니다: {uploaded_file.name}")
            # 문서 컨텍스트를 포함한 시스템 메시지 생성
            full_system_prompt = f"{system_prompt}\n\n참고할 문서 내용:\n\n{document_context}"
            st.session_state.initial_system_message = full_system_prompt
            
            # Nova Pro의 경우 초기 시스템 메시지를 user 메시지로 추가
            if model_name == "Nova Pro 1.0" and not st.session_state.nova_messages:
                st.session_state.nova_messages = [
                    ChatMessage('user', full_system_prompt),
                    ChatMessage('assistant', "네, 이해했습니다. 업로드된 문서를 참고하여 답변하도록 하겠습니다.")
                ]

    # 시스템 프롬프트 설정
    if st.session_state.initial_system_message and model_name == "Claude 3.7 Sonnet":
        system_prompt = st.session_state.initial_system_message

    conv_chain = init_conversation_chain(
        temperature, top_p, top_k, max_tokens, system_prompt, model_name, extended_thinking
    )

    # 저장된 메시지 표시
    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])

    # 사용자 입력 처리
    if prompt := st.chat_input("메시지를 입력하세요"):
        st.session_state.messages.append({"role": "user", "content": prompt})
        st.session_state.chat_history.add_user_message(prompt)
        
        with st.chat_message("user"):
            st.markdown(prompt)

        # 응답 생성
        response = generate_response(conv_chain, prompt, st.session_state.chat_history, show_reasoning)
        
        st.session_state.messages.append({"role": "assistant", "content": response})
        st.session_state.chat_history.add_ai_message(response)

if __name__ == "__main__":
    main()
