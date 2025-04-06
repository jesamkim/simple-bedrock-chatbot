import random
import json
from typing import List, Tuple, Union, Dict
import os
import tempfile
import csv
import pandas as pd
import warnings

# Pydantic 경고 무시
warnings.filterwarnings("ignore", category=UserWarning, module="pydantic")

import streamlit as st
from langchain_community.chat_message_histories import StreamlitChatMessageHistory
from langchain_aws import ChatBedrock
from langchain_core.messages import AIMessage, HumanMessage, SystemMessage
from langchain.callbacks.base import BaseCallbackHandler
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder
import boto3

# MCP 기능 임포트
import re
import json
from mcp_client import UnifiedMCPClient

# 통합 MCP 클라이언트 초기화
mcp_client = UnifiedMCPClient()
extract_keywords = mcp_client.extract_keywords

# 질의 의도 타입
class QueryIntent:
    DATETIME = "datetime"     # 날짜/시간 관련 질의
    SEARCH = "search"         # 검색이 필요한 질의
    GENERAL = "general"       # 특정 패턴이 없는 일반 질의
    MIXED = "mixed"           # 날짜/시간과 검색 모두 필요한 복합 의도
    TIME_COMPARISON = "time_comparison"  # 시간 비교 (기간, 경과 시간 등)

# 패턴 정의
# 시간 관련 패턴 - 명확한 시간 질의에만 매치되도록 구체화
TIME_PATTERNS = [
    r'\b(지금|현재|오늘)?\s*(몇\s*시|시간|시계)\b',
    r'\b(what|current)\s*time\b',
    r'\btime\s*(now|is it)\b'
]

# 날짜 관련 패턴 - 명확한 날짜 질의에만 매치되도록 구체화
DATE_PATTERNS = [
    r'\b(오늘|지금|현재)\s*(무슨|며칠|몇\s*일|날짜)\b',
    r'\b(무슨|며칠|몇\s*일|날짜)\b',
    r'\b(what|current)\s*(day|date)\b',
    r'\bdate\s*today\b',
    r'\btoday\s*is\b'
]

# 시간 비교 패턴 (두 시간/날짜 간의 관계 질의)
TIME_COMPARISON_PATTERNS = [
    r'\b(얼마나|며칠이|몇\s*일이|몇\s*년이|몇\s*개월이|시간이)\s*(지났|경과|됐|흘렀|남았|차이|벌어)\w*\b',
    r'\b(부터|이후|로부터|전부터|이전부터|이래|까지)\s*(얼마나|몇|지난)\w*\b',
    r'\b(몇\s*일\s*전|몇\s*일\s*후|몇\s*일\s*뒤)\b',
    r'\b(기간|간격|날짜\s*차이|시간\s*차이)\b',
    r'\b(since|elapsed|passed|ago|difference)\b'
]

# 현재 시간/날짜 관련 패턴 (날짜/시간 계산 시 현재 시점 지시)
CURRENT_TIME_PATTERNS = [
    r'\b(오늘|지금|현재|이제|now|today|current)\b'
]

# 주의해야 할 혼동 가능 단어 (이들이 포함되면 날짜/시간 질의가 아닐 가능성 높음)
AMBIGUOUS_TERMS = [
    r'\b(출시일|생일|기념일|공휴일|발매일|등록일|시작일|만기일)\b',
    r'\b(생년월일|설립일|계약일|입사일|퇴사일)\b'
]

# 날짜/시간 복합 패턴 (좀 더 구체적인 패턴만 포함)
DATETIME_PATTERNS = [
    r'\b지금|현재\s*(시간|날짜)\b',
    r'\bnow\b'
]

# 검색 의도를 나타내는 패턴
SEARCH_PATTERNS = [
    r'\b(뭐|무엇|어떤|어디|언제|왜|누구|어떻게)\b',
    r'\b(알려|찾아|검색|의미|뜻|방법|가격|위치|누구의|어디의)\b',
    r'\b(정보|사용법|차이|비교|종류|문제|이유|발매|사양)\b',
    r'\b(what|where|when|why|who|how)\b',
    r'\b(information|details|compare|difference|release|reason)\b'
]

REGION = os.getenv("AWS_REGION", "us-west-2")
if not REGION:
    raise ValueError("AWS_REGION 환경 변수가 설정되지 않았습니다.")

MODELS = {
    "Claude 3.7 Sonnet": {
        "id": "us.anthropic.claude-3-7-sonnet-20250219-v1:0",
        "class": ChatBedrock,
        "use_model_kwargs": True,
        "max_tokens": 8192
    }
}

SUPPORTED_FORMATS = ['pdf', 'doc', 'docx', 'md', 'ppt', 'pptx', 'txt', 'html', 'csv', 'xls', 'xlsx']
MAX_MESSAGES = 10  # 대화 기록 최대 유지 수

class ChatMessage:
    def __init__(self, role: str, text: str):
        self.role = role
        self.text = text

def initialize_session_state():
    """세션 상태 초기화"""
    if "widget_key" not in st.session_state:
        st.session_state["widget_key"] = str(random.randint(1, 1000000))
    if "messages" not in st.session_state:
        st.session_state.messages = []
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = StreamlitChatMessageHistory(key="chat_history")
    if "context" not in st.session_state:
        st.session_state.context = None
    if "initial_system_message" not in st.session_state:
        st.session_state.initial_system_message = None

class StreamHandler(BaseCallbackHandler):
    def __init__(self, container: st.container) -> None:
        self.container = container
        self.text = ""

    def on_llm_new_token(self, token: str, **kwargs) -> None:
        self.text += token
        self.container.markdown(self.text)

def set_page_config() -> None:
    st.set_page_config(page_title="Bedrock Chatbot", layout="wide")
    st.title("Bedrock Chatbot with Document Q&A")

def get_sidebar_params() -> Tuple[float, float, int, int, int, str, object, str, bool, bool, bool]:
    with st.sidebar:
        st.markdown("## 모델 정보")
        # 모델 선택 대신 고정 텍스트로 표시
        st.info("Claude 3.7 Sonnet v1")
        
        # 모델 이름은 항상 고정
        model_name = "Claude 3.7 Sonnet"
        
        # 모드 선택 (기본, MCP, Reasoning 중 하나만 선택 가능)
        mode = st.radio(
            "작동 모드 선택",
            options=["기본 모드", "MCP 모드", "Reasoning 모드"],
            index=1,
            help="기본 모드: 일반 챗봇 기능, MCP 모드: 웹 검색 및 현재 시간/날짜 정보 제공, Reasoning 모드: 복잡한 문제 해결에 특화된 사고 과정 제공",
            key=f"{st.session_state['widget_key']}_Mode"
        )
        
        # 선택된 모드에 따라 플래그 설정
        extended_thinking = False
        show_reasoning = False
        mcp_enable = False
        
        if mode == "Reasoning 모드":
            extended_thinking = True
            show_reasoning = True
            st.info("Reasoning 모드가 활성화되어 Temperature 값이 1.0으로 자동 설정되고 Top-K 및 Top-P 설정이 비활성화됩니다.")
            
            # Reasoning 모드에서만 표시 여부 옵션 제공
            show_reasoning = st.checkbox(
                "Reasoning 과정 표시",
                value=True,
                help="Claude의 사고 과정을 실시간으로 표시합니다.",
                key=f"{st.session_state['widget_key']}_Show_Reasoning"
            )
        elif mode == "MCP 모드":
            mcp_enable = True
            st.info("MCP 모드에서는 웹 검색 및 시간/날짜 정보를 자동으로 제공합니다.")
        
        st.markdown("## Document Upload")
        uploaded_file = st.file_uploader(
            "Upload a document",
            type=SUPPORTED_FORMATS,
            help="지원되는 파일 형식: pdf, doc, docx, md, ppt, pptx, txt, html, csv, xls, xlsx",
            key=f"{st.session_state['widget_key']}_file_uploader"  # widget_key와 연동
        )

        st.markdown("## Inference Parameters")
        system_prompt = st.text_area(
            "System Prompt",
            "You're a helpful assistant who loves to respond in Korean. You'll answer questions based on the uploaded documents, when a document exists. When you're not certain about something, don't guess or make up information. Instead, honestly admit that you don't know.",
            key=f"{st.session_state['widget_key']}_System_Prompt",
        )

        temperature = st.slider(
            "Temperature",
            min_value=0.0,
            max_value=1.0,
            value=0.0,
            step=0.1,
            key=f"{st.session_state['widget_key']}_Temperature",
        )

        with st.container():
            col1, col2 = st.columns(2)
            with col1:
                top_p = st.slider(
                    "Top-P",
                    min_value=0.0,
                    max_value=1.0,
                    value=1.00,
                    step=0.01,
                    key=f"{st.session_state['widget_key']}_Top-P",
                )
            with col2:
                top_k = st.slider(
                    "Top-K",
                    min_value=1,
                    max_value=500,
                    value=500,
                    step=5,
                    key=f"{st.session_state['widget_key']}_Top-K",
                )

        with st.container():
            col1, col2 = st.columns(2)
            with col1:
                model_max_tokens = MODELS[model_name]["max_tokens"]
                max_tokens = st.slider(
                    "Max Token",
                    min_value=0,
                    max_value=model_max_tokens,
                    value=model_max_tokens,
                    step=8,
                    key=f"{st.session_state['widget_key']}_Max_Token",
                )
            with col2:
                memory_window = st.slider(
                    "Memory Window",
                    min_value=0,
                    max_value=10,
                    value=10,
                    step=1,
                    key=f"{st.session_state['widget_key']}_Memory_Window",
                )

    return temperature, top_p, top_k, max_tokens, memory_window, system_prompt, uploaded_file, model_name, extended_thinking, show_reasoning, mcp_enable

def process_uploaded_file(file_path: str) -> str:
    """문서 파일을 처리하여 텍스트로 변환"""
    file_ext = os.path.splitext(file_path)[1].lower()
    
    try:
        if file_ext == '.pdf':
            from PyPDF2 import PdfReader
            reader = PdfReader(file_path)
            # PDF 파일이 암호화되어 있는 경우 처리되지 않음
            if reader.is_encrypted:
                raise ValueError("암호화된 PDF 파일은 처리할 수 없습니다.")
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            return text
        
        elif file_ext in ['.doc', '.docx']:
            from docx import Document
            doc = Document(file_path)
            text = ""
            for para in doc.paragraphs:
                text += para.text + "\n"
            return text
        
        elif file_ext in ['.txt', '.md', '.html']:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        
        elif file_ext in ['.csv']:
            df = pd.read_csv(file_path)
            return df.to_string()
        
        elif file_ext in ['.xls', '.xlsx']:
            df = pd.read_excel(file_path)
            return df.to_string()
        
        elif file_ext in ['.ppt', '.pptx']:
            from pptx import Presentation
            text = []
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n\n".join(text)
        
        else:
            raise ValueError(f"지원하지 않는 파일 형식입니다: {file_ext}")
            
    except Exception as e:
        raise Exception(f"파일 처리 중 오류가 발생했습니다 ({file_ext}): {str(e)}")

def init_conversation_chain(
    temperature: float,
    top_p: float,
    top_k: int,
    max_tokens: int,
    system_prompt: str,
    model_name: str,
    extended_thinking: bool = False,
) -> Tuple[boto3.client, dict]:

    model_info = MODELS[model_name]
    model_id = model_info["id"]
    
    # 직접 boto3 클라이언트 사용
    bedrock_client = boto3.client("bedrock-runtime", region_name=REGION)
    
    model_params = {
        "temperature": temperature,
        "top_p": top_p,
        "top_k": top_k,
        "max_tokens": max_tokens,
        "system": system_prompt,
        "model_id": model_id,
    }
    
    # Model reasoning 모드가 활성화된 경우
    if extended_thinking:
        # Model reasoning 모드에서는 temperature가 반드시 1이어야 함
        model_params["temperature"] = 1.0
        model_params["anthropic_version"] = "bedrock-2023-05-31"
        
        # Model reasoning 모드에서는 top_k와 top_p를 설정하지 않아야 함
        if "top_k" in model_params:
            del model_params["top_k"]
        if "top_p" in model_params:
            del model_params["top_p"]
        
        # 최대 Length는 64000으로 설정
        model_params["max_tokens"] = 64000
        
        # thinking.budget_tokens는 max_tokens보다 작아야 하며, 최대 4096
        thinking_budget = min(4096, model_params["max_tokens"] - 1000)
        
        model_params["thinking"] = {
            "type": "enabled",
            "budget_tokens": thinking_budget  # 사고 과정에 할당할 토큰 수
        }
    
    return bedrock_client, model_params

def convert_langchain_messages_to_anthropic(messages):
    """LangChain 메시지를 Anthropic API 형식으로 변환"""
    result = []
    system_message = None
    
    for msg in messages:
        if msg.type == "system":
            system_message = msg.content  # 시스템 메시지 별도 저장
        elif msg.type == "ai":
            result.append({"role": "assistant", "content": msg.content})
        elif msg.type == "human":
            result.append({"role": "user", "content": msg.content})
    
    return result, system_message

def convert_chat_messages_to_langchain(chat_messages: List[ChatMessage]) -> List[Union[HumanMessage, AIMessage, SystemMessage]]:
    messages = []
    for chat_msg in chat_messages:
        if not chat_msg:
            continue
        if chat_msg.role == "user":
            messages.append(HumanMessage(content=chat_msg.text))
        elif chat_msg.role == "assistant":
            messages.append(AIMessage(content=chat_msg.text))
        elif chat_msg.role == "system":
            messages.append(SystemMessage(content=chat_msg.text))
    return messages

def generate_response(
    conversation_data: Tuple[boto3.client, dict],
    input_text: str,
    chat_history: StreamlitChatMessageHistory,
    show_reasoning: bool = False,
    mcp_enable: bool = False
) -> str:
    # 입력 텍스트 길이 제한 체크
    if len(input_text) > 32000:  # Claude 3의 최대 입력 토큰 제한
        raise ValueError("입력 텍스트가 너무 깁니다.")

    # boto3 클라이언트와 모델 파라미터 추출
    client, model_params = conversation_data
    model_id = model_params["model_id"]
    system_message = model_params.get("system", "")
    
    with st.chat_message("assistant"):
        message_placeholder = st.empty()
        full_response = ""
        
        # MCP 서비스 처리 결과를 저장할 변수
        datetime_info_text = ""
        search_results_text = ""
        
        # MCP 활성화 상태에서 처리
        if mcp_enable:
            # 질의 분석 및 서비스 실행
            try:
                # 이전 메시지가 있으면 대화 컨텍스트 생성
                chat_context = None
                if chat_history.messages and len(chat_history.messages) >= 2:
                    # 가장 최근 메시지 2개(사용자 질문-모델 응답)
                    last_messages = chat_history.messages[-2:]
                    chat_context = [msg.content for msg in last_messages]
                
                # 1. Claude 3.7을 활용한 질의 의도 분석 (Thinking API)
                st.info("🧠 Claude로 질의 의도 분석 중...")
                
                # 대화 내용 추출 (최근 메시지 최대 4개)
                chat_context_list = None
                if chat_history and len(chat_history.messages) >= 2:
                    # 가장 최근 메시지 최대 4개 (2쌍)
                    messages = chat_history.messages[-4:] if len(chat_history.messages) >= 4 else chat_history.messages
                    chat_context_list = [msg.content for msg in messages]
                
                # LLM으로 의도 분석 (이전 대화 컨텍스트를 고려)
                intent_analysis = analyze_query_intent_with_llm(input_text, client, chat_context_list)
                
                # 분석 결과 표시
                intent = intent_analysis.get("intent", "general")
                subtype = intent_analysis.get("subtype", "none")
                datetime_needed = intent_analysis.get("datetime_needed", False)
                search_needed = intent_analysis.get("search_needed", True)
                reasoning = intent_analysis.get("reasoning", "")
                
                # 분석 결과 표시
                with st.expander("🧩 질의 의도 분석 결과", expanded=False):
                    st.markdown(f"""
                    **의도 유형:** {intent} {f'({subtype})' if subtype and subtype != 'none' else ''}
                    **날짜/시간 정보 필요:** {'✅' if datetime_needed else '❌'}
                    **웹 검색 필요:** {'✅' if search_needed else '❌'}
                    **분석 이유:** {reasoning}
                    """)
                
                # 2. 분석 결과에 따른 MCP 서비스 실행
                # 날짜/시간 정보가 필요한 경우
                if datetime_needed:
                    try:
                        st.info("⏰ 날짜/시간 정보를 조회합니다")
                        dt_info = mcp_client.get_datetime_info()
                        datetime_info_text = mcp_client.format_datetime_info(dt_info)
                        st.success("현재 날짜/시간 정보 조회 완료")
                        
                        # 결과 표시
                        with st.expander("📅 날짜/시간 정보"):
                            st.markdown(datetime_info_text)
                    except Exception as e:
                        st.error(f"날짜/시간 정보 조회 중 오류 발생: {str(e)}")
                
                # 검색 정보가 필요한 경우
                if search_needed:
                    # 검색 키워드 추출
                    keywords = extract_keywords(input_text)
                    search_query = " ".join(keywords)
                    
                    if search_query:
                        st.info(f"🔍 Google에서 '{search_query}'에 대한 정보 검색 중")
                        search_results = mcp_client.search(search_query)
                        
                        if search_results:
                            search_results_text = mcp_client.format_results(search_results)
                            st.success(f"'{search_query}' 검색 결과 {len(search_results)}건 발견")
                            
                            # 검색 결과 표시
                            with st.expander("🌐 검색 결과"):
                                st.markdown(search_results_text)
                        else:
                            st.warning(f"'{search_query}' 관련 검색 결과를 찾을 수 없습니다")
                
            except Exception as e:
                st.error(f"MCP 서비스 처리 중 오류 발생: {str(e)}")
        
        try:
            # LangChain 메시지 형식으로 변환
            messages = []
            messages.append(SystemMessage(content=system_message))
            for msg in chat_history.messages:
                messages.append(msg)
            
            # MCP 정보를 프롬프트에 포함
            if mcp_enable:
                # 검색 결과와 날짜/시간 정보가 모두 있는 경우
                if search_results_text and datetime_info_text:
                    enhanced_input = f"""질문: {input_text}
                    
다음은 질문과 관련된 정보입니다:

[인터넷 검색 결과]
{search_results_text}

[날짜/시간 정보]
{datetime_info_text}

위 정보를 참고하여 질문에 답변해주세요. 답변에서 인용한 정보가 있다면 반드시 그 출처를 명시해주세요."""
                    messages.append(HumanMessage(content=enhanced_input))
                
                # 검색 결과만 있는 경우
                elif search_results_text:
                    enhanced_input = f"""질문: {input_text}
                    
다음은 질문과 관련된 인터넷 검색 결과입니다:
{search_results_text}

위 정보를 참고하여 질문에 답변해주세요. 답변에서 인용한 정보가 있다면 반드시 그 출처(번호)를 [1], [2]와 같은 형식으로 명시해주세요."""
                    messages.append(HumanMessage(content=enhanced_input))
                
                # 날짜/시간 정보만 있는 경우
                elif datetime_info_text:
                    enhanced_input = f"""질문: {input_text}
                    
다음은 현재 날짜/시간 정보입니다:
{datetime_info_text}

위 정보를 참고하여 질문에 답변해주세요."""
                    messages.append(HumanMessage(content=enhanced_input))
                
                # 둘 다 없는 경우
                else:
                    messages.append(HumanMessage(content=input_text))
            else:
                # MCP가 비활성화된 경우 원본 질문만 전달
                messages.append(HumanMessage(content=input_text))
            
            # 메시지를 Claude API 형식으로 변환
            anthropic_messages, system_content = convert_langchain_messages_to_anthropic(messages)
            
            # MCP를 사용하는 경우 reasoning 모드는 비활성화
            if mcp_enable and "thinking" in model_params:
                st.warning("MCP 사용 시 Model reasoning 모드는 비활성화됩니다.")
                has_thinking = False
            else:
                has_thinking = "thinking" in model_params

            # API 요청 페이로드 구성
            request_payload = {
                "anthropic_version": model_params.get("anthropic_version", "bedrock-2023-05-31"),
                "max_tokens": model_params.get("max_tokens", 8192),
                "temperature": model_params.get("temperature", 0.0),
                "messages": anthropic_messages,
                "system": system_content
            }
            
            # reasoning 모드가 활성화된 경우 thinking 파라미터 추가
            if has_thinking and not mcp_enable:
                request_payload["thinking"] = model_params["thinking"]
                st.info("Model reasoning 모드로 응답 생성 중...")
            
            # top_p와 top_k 값이 있는 경우에만 추가
            if "top_p" in model_params:
                request_payload["top_p"] = model_params["top_p"]
            if "top_k" in model_params:
                request_payload["top_k"] = model_params["top_k"]
            
            try:
                # 스트리밍 응답 처리
                response = client.invoke_model_with_response_stream(
                    modelId=model_id,
                    body=json.dumps(request_payload)
                )
                
                # reasoning_placeholder는 reasoning 모드가 활성화된 경우에만 사용
                reasoning_placeholder = st.empty() if has_thinking and show_reasoning else None
                reasoning_text = ""
                
                for event in response["body"]:
                    try:
                        chunk = json.loads(event["chunk"]["bytes"])
                        
                        # thinking 타입 처리 (reasoning 과정)
                        if has_thinking and chunk.get("type") == "thinking":
                            if show_reasoning:
                                thinking_content = chunk.get("thinking", "")
                                reasoning_text += thinking_content
                                
                                # reasoning 과정을 UI에 표시
                                bg_color = "#2a2a2a"
                                text_color = "#ffffff"
                                
                                reasoning_placeholder.markdown(f"""
                                <div style="background-color: {bg_color}; color: {text_color}; padding: 10px; border-radius: 5px; margin-bottom: 10px; border-left: 4px solid #007bff; border: 1px solid #007bff; max-height: 400px; overflow-y: auto;">
                                    <h4 style="color: {text_color}; margin-top: 0;">🧠 Reasoning...</h4>
                                    <pre style="white-space: pre-wrap; overflow-wrap: break-word; color: {text_color}; background-color: transparent; border: none; padding: 0; margin: 0;">{reasoning_text}</pre>
                                </div>
                                """, unsafe_allow_html=True)
                        
                        # content_block_delta 타입 처리
                        elif chunk.get("type") == "content_block_delta":
                            if chunk["delta"].get("type") == "text" or chunk["delta"].get("type") == "text_delta":
                                text_chunk = chunk["delta"].get("text", "")
                                full_response += text_chunk
                                
                                # 일정 간격으로 UI 업데이트
                                if len(text_chunk) > 10 or text_chunk.endswith(('.', '!', '?', '\n')):
                                    message_placeholder.markdown(full_response + "▌")
                                    
                            # thinking_delta 처리
                            elif has_thinking and show_reasoning and chunk["delta"].get("type") == "thinking_delta":
                                thinking_chunk = chunk["delta"].get("thinking", "")
                                if thinking_chunk:
                                    reasoning_text += thinking_chunk
                                    
                                    # UI 업데이트
                                    bg_color = "#2a2a2a" 
                                    text_color = "#ffffff"
                                    
                                    reasoning_placeholder.markdown(f"""
                                    <div style="background-color: {bg_color}; color: {text_color}; padding: 10px; border-radius: 5px; margin-bottom: 10px; border-left: 4px solid #007bff; border: 1px solid #007bff; max-height: 400px; overflow-y: auto;">
                                        <h4 style="color: {text_color}; margin-top: 0;">🧠 Reasoning...</h4>
                                        <pre style="white-space: pre-wrap; overflow-wrap: break-word; color: {text_color}; background-color: transparent; border: none; padding: 0; margin: 0;">{reasoning_text}</pre>
                                    </div>
                                    """, unsafe_allow_html=True)
                        
                        # 응답 완료 처리
                        elif chunk.get("type") == "message_stop":
                            message_placeholder.markdown(full_response)
                    except Exception as chunk_error:
                        print(f"청크 처리 오류: {str(chunk_error)}")
                
                message_placeholder.markdown(full_response)
                return full_response
                
            except Exception as e:
                st.error(f"스트리밍 처리 중 오류: {str(e)}")
                # 스트리밍 실패 시 일반 호출 시도
                try:
                    response = client.invoke_model(
                        modelId=model_id,
                        body=json.dumps(request_payload)
                    )
                    response_body = json.loads(response["body"].read().decode("utf-8"))
                    full_response = response_body.get("completion", "") or response_body.get("content", "")
                    message_placeholder.markdown(full_response)
                    return full_response
                except Exception as e2:
                    st.error(f"일반 호출 처리 중 오류: {str(e2)}")
                    return f"응답 생성 중 오류가 발생했습니다: {str(e2)}"
            
        except Exception as e:
            error_detail = str(e)
            print(f"상세 오류: {error_detail}")  # 터미널에 자세한 오류 출력
            
            if "ThrottlingException" in error_detail:
                error_message = "요청을 처리하지 못했습니다. 잠시 후 다시 말씀해 주세요. 🙏"
            elif "validationException" in error_detail:
                error_message = "API 검증 오류가 발생했습니다. 요청 형식을 확인해주세요."
                print(f"API 검증 오류 상세: {error_detail}")
            else:
                error_message = f"죄송합니다. 오류가 발생했습니다: {error_detail}"
            
            message_placeholder.markdown(error_message)
            return error_message

def analyze_query_intent_with_llm(query: str, client, chat_history=None) -> Dict:
    """
    Claude 3.7 모델을 사용하여 질의 의도를 분석합니다.
    
    Args:
        query: 사용자 질의
        client: boto3 bedrock-runtime 클라이언트
        chat_history: 이전 대화 기록 (선택적)
        
    Returns:
        Dict: {
            "intent": 의도 (datetime, search, mixed, general 중 하나),
            "subtype": 세부 유형 (datetime인 경우 time, date 등),
            "datetime_needed": 날짜/시간 정보 필요 여부 (bool),
            "search_needed": 검색 필요 여부 (bool),
            "reasoning": 분석 이유
        }
    """
    # 이전 대화 컨텍스트 구성
    context = ""
    if chat_history and len(chat_history) > 0:
        last_messages = chat_history[-4:] if len(chat_history) > 4 else chat_history
        context = "이전 대화:\n" + "\n".join([f"{'사용자' if i % 2 == 0 else '어시스턴트'}: {msg}" for i, msg in enumerate(last_messages)])
    
    # 의도 분석을 위한 프롬프트
    system_prompt = """당신은 사용자 질의의 의도를 정확하게 분석하는 전문가입니다. 
사용자 질문을 분석하여 필요한 정보 유형과 서비스를 판단해주세요.

다음 두 가지 서비스를 활용할 수 있습니다:
1. 날짜/시간 정보: 현재 시간, 날짜, 요일 등의 정보
2. 웹 검색: 인터넷에서 특정 정보를 검색

분석해야 할 포인트:
- 현재 날짜/시간 정보가 필요한가? (예: "지금 몇 시야", "오늘은 무슨 요일이지")
- 웹 검색이 필요한가? (예: "마비노기 모바일 정보", "파이썬이란?")
- 두 서비스가 모두 필요한가? (예: "마비노기 모바일 출시일로부터 얼마나 지났어?")
  => 이런 질문은 "출시일"을 검색하고, "현재 날짜"를 알아야 답변 가능
  
특히 다음 경우들에 주의하세요:
* "~일로부터 얼마나 지났어?" 형태의 질문: 해당 일자 정보(검색)와 현재 날짜(날짜/시간 API) 두 가지 모두 필요함
* "출시일", "생일", "개업일", "설립일" 등을 포함한 질문: 검색이 필요한 경우가 많음
* 단순 시간/날짜 질의: "지금 몇 시야?", "오늘 날짜가 어떻게 돼?" 등은 날짜/시간 서비스만 필요

JSON 형식으로 아래와 같이 분석 결과를 반환해주세요:
{
  "intent": "datetime|search|mixed|general", 
  "subtype": "time|date|datetime|none",
  "datetime_needed": true/false,
  "search_needed": true/false,
  "reasoning": "분석에 대한 설명"
}
"""

    # 질의 텍스트
    query_text = f"""사용자 질문: {query}

{context}

이 질문에 답변하기 위해 필요한 서비스를 분석해주세요."""

    try:
        # Claude 3.7에 메시지 전송
        response = client.invoke_model(
            modelId="us.anthropic.claude-3-7-sonnet-20250219-v1:0",
            body=json.dumps({
                "anthropic_version": "bedrock-2023-05-31",
                "max_tokens": 1000,
                "temperature": 0,
                "system": system_prompt,
                "messages": [
                    {"role": "user", "content": query_text}
                ]
            })
        )
        
        # 응답 파싱
        response_body = json.loads(response['body'].read())
        response_content = response_body.get('content', [{'text': '{"intent": "general", "datetime_needed": false, "search_needed": false, "reasoning": "분석 실패"}'}])[0]['text']
        
        # JSON 추출 (응답 텍스트에 JSON 이외의 내용이 있을 수 있으므로)
        json_start = response_content.find('{')
        json_end = response_content.rfind('}') + 1
        if json_start >= 0 and json_end > json_start:
            json_str = response_content[json_start:json_end]
            try:
                result = json.loads(json_str)
                
                # 필수 키가 없는 경우 기본값 추가
                if "intent" not in result:
                    result["intent"] = "general"
                if "datetime_needed" not in result:
                    result["datetime_needed"] = False
                if "search_needed" not in result:
                    result["search_needed"] = False
                if "subtype" not in result:
                    result["subtype"] = "none"
                
                return result
            except json.JSONDecodeError:
                print(f"JSON 파싱 오류: {json_str}")
                # 기본값 반환
                return {
                    "intent": "general",
                    "subtype": "none",
                    "datetime_needed": False,
                    "search_needed": True,
                    "reasoning": "의도 분석 중 오류 발생"
                }
        else:
            print(f"JSON을 찾을 수 없음: {response_content}")
            # 기본값 반환
            return {
                "intent": "general", 
                "subtype": "none",
                "datetime_needed": False, 
                "search_needed": True,
                "reasoning": "의도 분석 실패"
            }
            
    except Exception as e:
        print(f"의도 분석 요청 중 오류: {e}")
        # 오류 발생 시 기본값 반환
        return {
            "intent": "general",
            "subtype": "none", 
            "datetime_needed": False,
            "search_needed": True,
            "reasoning": f"오류: {str(e)}"
        }

def analyze_query_intent(query: str, chat_history=None) -> tuple[str, str]:
    """
    사용자 질의의 의도를 분석하여 서비스 유형과 세부 유형을 반환합니다.
    규칙 기반 패턴 매칭 방식 (레거시 버전)
    
    Args:
        query: 사용자 질의
        chat_history: 이전 대화 기록 (선택적)
        
    Returns:
        tuple: (의도 유형, 세부 유형)
            의도 유형: "datetime", "search", "mixed", "general" 중 하나
            세부 유형: datetime인 경우 "time", "date", "datetime" 중 하나
    """
    query_lower = query.lower()
    
    # 각 패턴 매칭 결과 저장
    pattern_matches = {
        "time": any(re.search(pattern, query_lower) for pattern in TIME_PATTERNS),
        "date": any(re.search(pattern, query_lower) for pattern in DATE_PATTERNS),
        "datetime": any(re.search(pattern, query_lower) for pattern in DATETIME_PATTERNS),
        "current_time": any(re.search(pattern, query_lower) for pattern in CURRENT_TIME_PATTERNS),
        "time_comparison": any(re.search(pattern, query_lower) for pattern in TIME_COMPARISON_PATTERNS),
        "search": any(re.search(pattern, query_lower) for pattern in SEARCH_PATTERNS),
        "ambiguous": any(re.search(pattern, query_lower) for pattern in AMBIGUOUS_TERMS)
    }
    
    # 1. 시간 비교 패턴 감지 (날짜 간 계산 등)
    if pattern_matches["time_comparison"]:
        # 시간 비교 + 현재 시간 언급 + 검색/임의어 = 복합 의도
        if (pattern_matches["current_time"] or pattern_matches["date"] or pattern_matches["time"]) and \
           (pattern_matches["search"] or pattern_matches["ambiguous"]):
            return QueryIntent.MIXED, "datetime"
        # 시간 비교 + 현재 시간 언급 = 날짜 시간 의도
        elif pattern_matches["current_time"] or pattern_matches["date"] or pattern_matches["time"]:
            return QueryIntent.DATETIME, "datetime"
    
    # 2. 혼동 가능어가 있지만, 시간 비교 패턴이 있는 경우 = 복합 의도
    if pattern_matches["ambiguous"] and (pattern_matches["time_comparison"] or (pattern_matches["current_time"] and pattern_matches["search"])):
        return QueryIntent.MIXED, "datetime"
    
    # 3. 날짜/시간 패턴과 검색 패턴이 모두 있는 경우 = 복합 의도
    if (pattern_matches["time"] or pattern_matches["date"] or pattern_matches["datetime"]) and \
       (pattern_matches["search"] or pattern_matches["ambiguous"]):
        return QueryIntent.MIXED, "datetime"
    
    # 4. 날짜/시간 패턴만 있는 경우
    if pattern_matches["time"] or pattern_matches["date"] or pattern_matches["datetime"]:
        # 시간과 날짜 패턴이 모두 있으면 datetime
        if pattern_matches["time"] and (pattern_matches["date"] or pattern_matches["datetime"]):
            return QueryIntent.DATETIME, "datetime"
        # 시간 패턴만 있으면 time
        elif pattern_matches["time"]:
            return QueryIntent.DATETIME, "time"
        # 날짜 패턴만 있으면 date
        else:
            return QueryIntent.DATETIME, "date"
    
    # 5. 검색 관련 패턴이 있는 경우 = 검색 의도
    if pattern_matches["search"] or pattern_matches["ambiguous"]:
        return QueryIntent.SEARCH, ""
    
    # 6. 특정 패턴이 감지되지 않으면 일반 의도로 간주
    return QueryIntent.GENERAL, ""

def process_mcp_services(input_text: str) -> tuple[str, str]:
    """
    사용자 입력을 분석하여 필요한 MCP 서비스를 결정하고 수행합니다.
    
    Args:
        input_text: 사용자 입력 텍스트
        
    Returns:
        tuple: (날짜/시간 정보, 검색 결과)
    """
    # 사용자 입력 분석
    intent, subtype = analyze_query_intent(input_text)
    datetime_info_text = ""
    search_results_text = ""
    
    st.info(f"MCP 서비스 분석 중: {intent} 의도 감지됨")
    
    # 의도에 따른 서비스 호출
    if intent == QueryIntent.DATETIME or intent == QueryIntent.MIXED:
        try:
            # 통합 MCP 클라이언트 사용
            if subtype == "time":
                # 시간 정보 요청
                time_info = mcp_client.get_current_time()
                datetime_info_text = mcp_client.format_time(time_info)
                st.success("시간 정보 조회 완료")
            elif subtype == "date":
                # 날짜 정보 요청
                date_info = mcp_client.get_current_date()
                datetime_info_text = mcp_client.format_date(date_info)
                st.success("날짜 정보 조회 완료")
            else:  # "datetime"
                # 종합 정보 요청
                dt_info = mcp_client.get_datetime_info()
                datetime_info_text = mcp_client.format_datetime_info(dt_info)
                st.success("날짜 및 시간 정보 조회 완료")
            
            # 결과 표시
            with st.expander("날짜/시간 정보 확인"):
                st.markdown(datetime_info_text)
                
        except Exception as e:
            st.error(f"날짜/시간 정보 조회 중 오류 발생: {str(e)}")
    
    if intent == QueryIntent.SEARCH or intent == QueryIntent.MIXED:
        try:
            # 키워드 추출
            keywords = extract_keywords(input_text)
            search_query = " ".join(keywords)
            
            if search_query:
                st.info(f"Google에서 '{search_query}'에 대한 정보를 검색 중입니다...")
                
                # 통합 MCP 클라이언트를 통한 검색 수행
                search_results = mcp_client.search(search_query)
                
                # 검색 결과 포맷팅
                search_results_text = mcp_client.format_results(search_results)
                
                if search_results:
                    st.success(f"'{search_query}' 관련 검색 완료")
                    
                    # 검색 결과를 보여주기 (접을 수 있게)
                    with st.expander("검색 결과 확인"):
                        st.markdown(search_results_text)
                else:
                    st.warning(f"'{search_query}' 관련 검색 결과가 없습니다.")
        except Exception as e:
            st.error(f"검색 중 오류가 발생했습니다: {str(e)}")
    
    return datetime_info_text, search_results_text

def new_chat() -> None:
    """새로운 대화 시작"""
    # widget_key를 새로 생성하여 file_uploader를 포함한 모든 위젯 초기화
    st.session_state["widget_key"] = str(random.randint(1, 1000000))
    st.session_state.messages = []
    st.session_state.chat_history.clear()
    st.session_state.context = None
    st.session_state.initial_system_message = None

def handle_file_upload(uploaded_file) -> str:
    """파일 업로드 처리"""
    if uploaded_file:
        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(uploaded_file.name)[1]) as tmp_file:
            tmp_file.write(uploaded_file.getvalue())
            tmp_file_path = tmp_file.name

        try:
            document_text = process_uploaded_file(tmp_file_path)
            st.session_state.context = document_text  # 대용량 문서의 경우 메모리 부족 발생 가능
            os.unlink(tmp_file_path)
            return document_text
        except Exception as e:
            os.unlink(tmp_file_path)  # 에러 발생 시에도 임시 파일 삭제 필요
            st.error(f"문서 처리 중 오류가 발생했습니다: {str(e)}")
            return None
    return None

def main() -> None:
    set_page_config()
    initialize_session_state()

    st.sidebar.button("New Chat", on_click=new_chat, type="primary")

    temperature, top_p, top_k, max_tokens, memory_window, system_prompt, uploaded_file, model_name, extended_thinking, show_reasoning, mcp_enable = get_sidebar_params()

    # 문서가 업로드되면 시스템 메시지 초기화
    if uploaded_file:
        document_context = handle_file_upload(uploaded_file)
        if document_context:
            st.sidebar.success(f"문서가 성공적으로 업로드되었습니다: {uploaded_file.name}")
            # 문서 컨텍스트를 포함한 시스템 메시지 생성
            full_system_prompt = f"{system_prompt}\n\n참고할 문서 내용:\n\n{document_context}"
            st.session_state.initial_system_message = full_system_prompt

    # 시스템 프롬프트 설정
    if st.session_state.initial_system_message:
        system_prompt = st.session_state.initial_system_message

    conv_chain = init_conversation_chain(
        temperature, top_p, top_k, max_tokens, system_prompt, model_name, extended_thinking
    )

    # 저장된 메시지 표시
    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])

    # 사용자 입력 처리
    if prompt := st.chat_input("메시지를 입력하세요"):
        st.session_state.messages.append({"role": "user", "content": prompt})
        st.session_state.chat_history.add_user_message(prompt)
        
        with st.chat_message("user"):
            st.markdown(prompt)

        # 응답 생성 및 세션 저장 (UI 표시는 generate_response에서 이미 처리됨)
        response = generate_response(conv_chain, prompt, st.session_state.chat_history, show_reasoning, mcp_enable)
        
        # 세션 상태 업데이트만 수행 (UI 표시는 하지 않음)
        st.session_state.messages.append({"role": "assistant", "content": response})
        st.session_state.chat_history.add_ai_message(response)

if __name__ == "__main__":
    main()
